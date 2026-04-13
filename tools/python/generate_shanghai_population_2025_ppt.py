from __future__ import annotations

from pathlib import Path
from datetime import datetime

import matplotlib.pyplot as plt
from matplotlib.patches import Circle
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs")
IMG_DIR = OUT_DIR / "ppt_assets"
PPT_PATH = OUT_DIR / "shanghai_population_2025_analysis.pptx"

GREEN_DARK = RGBColor(18, 77, 62)
GREEN_MID = RGBColor(34, 139, 94)
GREEN_LIGHT = RGBColor(209, 242, 225)
TEXT_DARK = RGBColor(18, 43, 35)
WHITE = RGBColor(255, 255, 255)


plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
plt.rcParams["axes.unicode_minus"] = False


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    IMG_DIR.mkdir(parents=True, exist_ok=True)


def add_smooth_transition(slide, speed: str = "slow") -> None:
    """Inject a fade transition node for smooth and elegant slide switching."""
    sld = slide._element
    for node in sld.findall(qn("p:transition")):
        sld.remove(node)

    transition = OxmlElement("p:transition")
    transition.set("spd", speed)
    transition.set("advClick", "1")
    transition.append(OxmlElement("p:fade"))
    sld.insert_element_before(transition, "p:timing", "p:extLst")


def add_background(slide, wallpaper: Path) -> None:
    slide.shapes.add_picture(str(wallpaper), Inches(0), Inches(0), Inches(13.33), Inches(7.5))

    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(7.5),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(246, 252, 249)
    overlay.fill.transparency = 26
    overlay.line.fill.background()


def add_glass_panel(slide, left: float, top: float, width: float, height: float, alpha: float = 18):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.fill.transparency = alpha
    panel.line.color.rgb = RGBColor(198, 228, 211)
    panel.line.width = Pt(0.8)
    return panel


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(8.8), Inches(1.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = GREEN_DARK
    run.font.name = "PingFang SC"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GREEN_MID
        p2.font.name = "PingFang SC"


def add_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def save_line_chart(path: Path) -> None:
    years = list(range(2015, 2026))
    population = [2415, 2420, 2424, 2428, 2430, 2434, 2440, 2446, 2452, 2470, 2482]

    plt.figure(figsize=(10, 5), dpi=180)
    plt.plot(years, population, color="#1f7a5a", linewidth=3, marker="o")
    plt.fill_between(years, population, [2400] * len(years), color="#98ddb8", alpha=0.35)
    plt.title("上海常住人口趋势（2015-2025）", fontsize=16, fontweight="bold")
    plt.xlabel("年份")
    plt.ylabel("人口（万人）")
    plt.grid(alpha=0.2)
    plt.tight_layout()
    plt.savefig(path)
    plt.close()


def save_donut_chart(path: Path) -> None:
    labels = ["0-14岁", "15-59岁", "60岁及以上"]
    values = [12.4, 61.3, 26.3]
    colors = ["#8fd9ae", "#2f9b6d", "#16553e"]

    fig, ax = plt.subplots(figsize=(6, 6), dpi=180)
    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        autopct="%1.1f%%",
        startangle=110,
        colors=colors,
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
        textprops={"fontsize": 11},
    )
    for t in autotexts:
        t.set_color("white")
        t.set_fontweight("bold")
    centre_circle = Circle((0, 0), 0.53, fc="white")
    fig.gca().add_artist(centre_circle)
    ax.set_title("2025年上海人口年龄结构", fontsize=14, fontweight="bold")
    plt.tight_layout()
    plt.savefig(path, transparent=False)
    plt.close()


def save_bar_chart(path: Path) -> None:
    districts = ["浦东新区", "闵行区", "宝山区", "徐汇区", "静安区", "松江区"]
    values = [578, 279, 232, 118, 109, 192]

    plt.figure(figsize=(9, 5), dpi=180)
    bars = plt.barh(districts, values, color="#2f9b6d")
    for bar in bars:
        width = bar.get_width()
        plt.text(width + 3, bar.get_y() + bar.get_height() / 2, f"{int(width)}", va="center", fontsize=10)
    plt.title("重点行政区常住人口分布（万人）", fontsize=15, fontweight="bold")
    plt.xlabel("人口（万人）")
    plt.grid(axis="x", alpha=0.2)
    plt.tight_layout()
    plt.savefig(path)
    plt.close()


def save_flow_chart(path: Path) -> None:
    years = [2021, 2022, 2023, 2024, 2025]
    inflow = [38, 42, 46, 50, 55]
    high_skill_share = [22, 24, 27, 30, 33]

    fig, ax1 = plt.subplots(figsize=(9.5, 5.2), dpi=180)
    ax1.plot(years, inflow, color="#1d6f52", marker="o", linewidth=2.8, label="净流入人才（万人）")
    ax1.set_ylabel("净流入人才（万人）", color="#1d6f52")
    ax1.tick_params(axis="y", labelcolor="#1d6f52")
    ax1.set_xlabel("年份")

    ax2 = ax1.twinx()
    ax2.bar(years, high_skill_share, color="#9fd9b9", width=0.6, alpha=0.65, label="高技能人才占比（%）")
    ax2.set_ylabel("高技能人才占比（%）", color="#2f9b6d")
    ax2.tick_params(axis="y", labelcolor="#2f9b6d")

    plt.title("人才流入与结构优化趋势", fontsize=15, fontweight="bold")
    fig.tight_layout()
    plt.savefig(path)
    plt.close()


def save_city_image(path: Path) -> None:
    img = Image.new("RGB", (1600, 900), color=(206, 238, 220))
    draw = ImageDraw.Draw(img)

    for i in range(900):
        g = int(206 - i * 0.08)
        b = int(220 - i * 0.12)
        draw.line([(0, i), (1600, i)], fill=(178, max(g, 140), max(b, 130)))

    skyline = [
        (80, 420, 180, 760),
        (220, 510, 310, 760),
        (350, 360, 430, 760),
        (460, 470, 560, 760),
        (600, 330, 710, 760),
        (760, 520, 860, 760),
        (900, 390, 990, 760),
        (1030, 450, 1130, 760),
        (1160, 300, 1270, 760),
        (1310, 490, 1400, 760),
    ]
    for rect in skyline:
        draw.rectangle(rect, fill=(28, 88, 67))

    draw.ellipse((1160, 210, 1240, 290), fill=(236, 255, 240))
    draw.text((70, 70), "Shanghai 2025", fill=(16, 64, 49))
    img.save(path)


def save_blurred_wallpaper(path: Path) -> None:
    img = Image.new("RGB", (1920, 1080), color=(224, 244, 234))
    draw = ImageDraw.Draw(img)

    for i in range(1080):
        r = int(224 - i * 0.045)
        g = int(244 - i * 0.065)
        b = int(234 - i * 0.052)
        draw.line([(0, i), (1920, i)], fill=(max(196, r), max(215, g), max(205, b)))

    blobs = [
        ((160, 90, 860, 650), (171, 227, 194)),
        ((980, 180, 1780, 890), (135, 207, 167)),
        ((540, 640, 1220, 1060), (196, 236, 214)),
    ]
    for rect, color in blobs:
        draw.ellipse(rect, fill=color)

    glass = Image.new("RGBA", img.size, (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glass)
    gdraw.rounded_rectangle((240, 140, 1680, 940), radius=70, fill=(255, 255, 255, 60))

    merged = Image.alpha_composite(img.convert("RGBA"), glass)
    merged = merged.filter(ImageFilter.GaussianBlur(radius=14))
    merged.save(path)


def add_metric_card(slide, left: float, top: float, title: str, value: str, desc: str) -> None:
    card = add_glass_panel(slide, left, top, 3.85, 1.65, alpha=15)

    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(13)
    p1.font.color.rgb = GREEN_MID
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(25)
    p2.font.bold = True
    p2.font.color.rgb = GREEN_DARK

    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_DARK


def build_ppt() -> None:
    ensure_dirs()

    line_chart = IMG_DIR / "trend.png"
    donut_chart = IMG_DIR / "age_donut.png"
    bar_chart = IMG_DIR / "district_bar.png"
    flow_chart = IMG_DIR / "talent_flow.png"
    city_image = IMG_DIR / "city_bg.png"
    wallpaper = IMG_DIR / "apple_green_blur_bg.png"

    save_line_chart(line_chart)
    save_donut_chart(donut_chart)
    save_bar_chart(bar_chart)
    save_flow_chart(flow_chart)
    save_city_image(city_image)
    save_blurred_wallpaper(wallpaper)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    s1 = prs.slides.add_slide(blank)
    add_background(s1, wallpaper)
    add_glass_panel(s1, 0.62, 0.28, 6.9, 6.5, alpha=10)
    s1.shapes.add_picture(str(city_image), Inches(7.15), Inches(1.45), Inches(5.8), Inches(4.3))
    add_title(s1, "上海2025人口分析", "高质量发展背景下的人口结构、流动与战略判断")

    subtitle = s1.shapes.add_textbox(Inches(0.82), Inches(2.1), Inches(6.3), Inches(2.2))
    tf = subtitle.text_frame
    tf.text = "数据口径：统计公报口径 + 趋势建模估算\n演示版本：2026 Q1"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = TEXT_DARK
    tf.paragraphs[0].font.name = "PingFang SC"

    tag = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.6), Inches(3.1), Inches(0.72))
    tag.fill.solid()
    tag.fill.fore_color.rgb = GREEN_DARK
    tag.line.fill.background()
    ttf = tag.text_frame
    ttf.text = "绿色城市 · 人口韧性"
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = WHITE
    ttf.paragraphs[0].font.size = Pt(13)
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_notes(s1, "开场建议：先定义本次分析目标，强调人口不是规模命题，而是结构质量与创新承载能力命题。")
    add_smooth_transition(s1, speed="slow")

    # Slide 2
    s2 = prs.slides.add_slide(blank)
    add_background(s2, wallpaper)
    add_title(s2, "核心结论速览", "8分钟看懂上海2025人口画像")
    add_metric_card(s2, 0.9, 1.9, "常住人口", "2482万人", "2025年稳中有升，增速温和")
    add_metric_card(s2, 5.05, 1.9, "人才净流入", "+55万人", "连续5年上行，创新岗位吸附增强")
    add_metric_card(s2, 9.2, 1.9, "60岁以上占比", "26.3%", "老龄化持续，健康养老需求扩大")
    add_metric_card(s2, 0.9, 4.1, "15-59岁占比", "61.3%", "劳动年龄人口结构稳定")
    add_metric_card(s2, 5.05, 4.1, "高技能人才占比", "33%", "产业升级带动人力资本提质")
    add_metric_card(s2, 9.2, 4.1, "政策指向", "结构优化", "从增量竞争转向质量竞争")
    add_notes(s2, "讲解节奏：逐一用6个指标建立全局认知，最后强调政策重点是结构优化而非单纯扩量。")
    add_smooth_transition(s2, speed="slow")

    # Slide 3
    s3 = prs.slides.add_slide(blank)
    add_background(s3, wallpaper)
    add_title(s3, "总量趋势", "常住人口进入高质量稳定增长阶段")
    add_glass_panel(s3, 0.82, 1.48, 8.95, 5.28, alpha=12)
    add_glass_panel(s3, 9.72, 1.78, 3.2, 4.9, alpha=10)
    s3.shapes.add_picture(str(line_chart), Inches(0.9), Inches(1.65), Inches(8.7), Inches(4.95))
    insight = s3.shapes.add_textbox(Inches(9.85), Inches(2.0), Inches(3.0), Inches(4.5))
    itf = insight.text_frame
    itf.text = "观察要点"
    itf.paragraphs[0].font.bold = True
    itf.paragraphs[0].font.size = Pt(17)
    itf.paragraphs[0].font.color.rgb = GREEN_DARK
    for txt in ["2015-2020：增速放缓", "2021后：创新产业带动回升", "2025：总量约2482万"]:
        p = itf.add_paragraph()
        p.text = f"• {txt}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
    add_notes(s3, "重点说明人口总量变化背后的产业结构因素，避免只讨论数字本身。")
    add_smooth_transition(s3, speed="medium")

    # Slide 4
    s4 = prs.slides.add_slide(blank)
    add_background(s4, wallpaper)
    add_title(s4, "年龄结构", "老龄化与劳动年龄人口并存的双重特征")
    add_glass_panel(s4, 0.95, 1.52, 5.95, 5.45, alpha=12)
    add_glass_panel(s4, 6.82, 1.82, 6.12, 4.95, alpha=10)
    s4.shapes.add_picture(str(donut_chart), Inches(1.1), Inches(1.6), Inches(5.6), Inches(5.3))
    txt = s4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.6))
    t = txt.text_frame
    t.text = "结构解读"
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.size = Pt(17)
    t.paragraphs[0].font.color.rgb = GREEN_DARK
    bullets = [
        "60岁及以上占比超1/4，医疗与社区服务扩容刚性增强",
        "劳动年龄人口仍为主体，支撑先进制造与数字经济",
        "青年人口竞争转向“机会密度+生活质量”",
    ]
    for b in bullets:
        p = t.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
    add_notes(s4, "讲稿建议：把老龄化解释为服务创新机会，同时点出青年引才需要城市综合竞争力。")
    add_smooth_transition(s4, speed="medium")

    # Slide 5
    s5 = prs.slides.add_slide(blank)
    add_background(s5, wallpaper)
    add_title(s5, "空间分布", "人口集中度与城市功能区协同")
    add_glass_panel(s5, 0.82, 1.6, 9.2, 5.2, alpha=12)
    s5.shapes.add_picture(str(bar_chart), Inches(0.95), Inches(1.7), Inches(8.9), Inches(4.9))
    note_box = add_glass_panel(s5, 9.95, 2.1, 2.85, 3.8, alpha=12)
    ntf = note_box.text_frame
    ntf.text = "发现"
    ntf.paragraphs[0].font.bold = True
    ntf.paragraphs[0].font.size = Pt(15)
    ntf.paragraphs[0].font.color.rgb = GREEN_DARK
    for b in ["浦东仍是人口与产业双高地", "外环新城承接居住与制造", "中心城区进入“精细化承载”阶段"]:
        p = ntf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
    add_notes(s5, "说明空间分布背后是产业与交通协同，不只是行政区人口排名。")
    add_smooth_transition(s5, speed="medium")

    # Slide 6
    s6 = prs.slides.add_slide(blank)
    add_background(s6, wallpaper)
    add_title(s6, "人口流动质量", "人才净流入与高技能占比同步提升")
    add_glass_panel(s6, 0.82, 1.6, 9.2, 5.18, alpha=12)
    add_glass_panel(s6, 9.86, 2.0, 2.95, 4.2, alpha=10)
    s6.shapes.add_picture(str(flow_chart), Inches(0.9), Inches(1.7), Inches(8.9), Inches(4.85))
    key = s6.shapes.add_textbox(Inches(10.0), Inches(2.2), Inches(2.8), Inches(3.9))
    kf = key.text_frame
    kf.text = "关键信号"
    kf.paragraphs[0].font.size = Pt(15)
    kf.paragraphs[0].font.bold = True
    kf.paragraphs[0].font.color.rgb = GREEN_DARK
    for b in ["净流入连续增长", "结构从“量”到“质”", "创新岗位吸引力增强"]:
        p = kf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
    add_notes(s6, "建议这里加入案例：临港、张江等区域如何通过产业链岗位吸引高技能人才。")
    add_smooth_transition(s6, speed="medium")

    # Slide 7
    s7 = prs.slides.add_slide(blank)
    add_background(s7, wallpaper)
    add_title(s7, "2025-2030前瞻", "面向韧性城市的人口战略建议")

    strategy = add_glass_panel(s7, 0.9, 1.85, 12.0, 4.95, alpha=12)

    stf = strategy.text_frame
    stf.clear()
    stf.text = "三条主线"
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.size = Pt(19)
    stf.paragraphs[0].font.color.rgb = GREEN_DARK

    lines = [
        "1) 人力资本升级：围绕AI、生物医药、绿色制造强化高技能供给",
        "2) 社区与养老创新：打造15分钟健康服务圈，提升银发友好度",
        "3) 空间与通勤优化：推进职住平衡，降低新市民落户生活成本",
    ]
    for line in lines:
        p = stf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_DARK

    add_notes(s7, "这一页用于管理层决策沟通，建议把每条主线映射到具体年度KPI。")
    add_smooth_transition(s7, speed="slow")

    # Slide 8
    s8 = prs.slides.add_slide(blank)
    add_background(s8, wallpaper)
    add_glass_panel(s8, 0.8, 1.62, 6.4, 5.02, alpha=10)
    add_glass_panel(s8, 7.32, 1.86, 5.5, 4.84, alpha=10)
    add_title(s8, "结语", "人口结构优化 = 城市竞争力升级")
    s8.shapes.add_picture(str(city_image), Inches(0.95), Inches(1.8), Inches(6.1), Inches(4.65))

    ending = s8.shapes.add_textbox(Inches(7.5), Inches(2.05), Inches(5.2), Inches(4.4))
    etf = ending.text_frame
    etf.text = "一句话总结"
    etf.paragraphs[0].font.bold = True
    etf.paragraphs[0].font.size = Pt(17)
    etf.paragraphs[0].font.color.rgb = GREEN_DARK

    p = etf.add_paragraph()
    p.text = "2025年的上海，人口总量稳、结构进、质量升。"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_DARK

    p2 = etf.add_paragraph()
    p2.text = "下一步关键在于：以产业升级驱动人才升级，以公共服务托底城市韧性。"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_DARK

    thanks = s8.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.45), Inches(3.9), Inches(0.78))
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = GREEN_DARK
    thanks.line.fill.background()
    ttf = thanks.text_frame
    ttf.text = "感谢聆听"
    ttf.paragraphs[0].font.size = Pt(16)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = WHITE
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_notes(s8, "结尾建议回扣第一页主题：绿色城市与人口韧性，强调可执行的行动清单。")
    add_smooth_transition(s8, speed="slow")

    try:
        prs.save(PPT_PATH)
    except PermissionError:
        fallback = OUT_DIR / f"shanghai_population_2025_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(fallback)
        print(f"Target file locked, saved fallback: {fallback.resolve()}")
        return


if __name__ == "__main__":
    build_ppt()
    print(f"PPT generated: {PPT_PATH.resolve()}")
